# -*- coding: utf-8 -*-
"""Generate a presentation PPTX summarizing the lane detection project,
including all 37 test-difficult result images."""
import json
import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

RESULTS_DIR = "results/test-difficult-clrnet"
OUT_PATH = "docs/道路检测_展示报告.pptx"
FONT = "Microsoft YaHei"

DARK = RGBColor(0x1F, 0x29, 0x37)
ACCENT = RGBColor(0x1E, 0x6F, 0xB8)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
GREY = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _set_font(run, size, bold=False, color=DARK):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT


def add_textbox(slide, left, top, width, height, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for index, (text, size, bold, color) in enumerate(lines):
        para = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        para.alignment = align
        run = para.add_run()
        run.text = text
        _set_font(run, size, bold, color)
    return box


def add_band(slide, top, height, color):
    shape = slide.shapes.add_shape(1, 0, top, SLIDE_W, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_band(slide, 0, SLIDE_H, DARK)
    add_textbox(slide, Inches(1), Inches(2.3), Inches(11.3), Inches(1.6),
                [("道路信息检测：车道线检测与黄/白颜色分类", 40, True, WHITE)])
    add_textbox(slide, Inches(1), Inches(4.0), Inches(11.3), Inches(1.8),
                [("CLRNet 检测 + 道路相对蓝色亏损颜色分类 + 轻量误检后处理", 22, False, RGBColor(0xCF, 0xE3, 0xF5)),
                 ("期末实验展示报告 · 测试集 test-difficult（37 张）", 18, False, RGBColor(0xAA, 0xC4, 0xDA))])


def section_title(slide, text):
    add_band(slide, 0, Inches(1.0), ACCENT)
    add_textbox(slide, Inches(0.6), Inches(0.18), Inches(12), Inches(0.7),
                [(text, 26, True, WHITE)], anchor=MSO_ANCHOR.MIDDLE)


def bullets_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    section_title(slide, title)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.7), Inches(5.6))
    tf = box.text_frame
    tf.word_wrap = True
    for index, (text, level, size, bold, color) in enumerate(bullets):
        para = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        para.level = level
        para.space_after = Pt(8)
        run = para.add_run()
        run.text = text
        _set_font(run, size, bold, color)
    return slide


def metrics_table_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    section_title(slide, "指标对比（对照老师标注真值，37 张）")
    rows = [
        ["模型版本", "白线 P / R", "黄线 P / R", "全部 P / R"],
        ["UFLD 原始颜色分类", "74.0% / 53.8%", "42.9% / 50.0%", "69.2% / 53.4%"],
        ["UFLD + 改进颜色分类", "72.8% / 55.7%", "80.0% / 66.7%", "73.6% / 56.8%"],
        ["CLRNet + 改进颜色分类", "84.4% / 76.4%", "75.0% / 75.0%", "83.3% / 76.3%"],
        ["CLRNet + 颜色分类 + 轻量误检过滤", "88.0% / 76.4%", "90.0% / 75.0%", "88.2% / 76.3%"],
    ]
    n_rows, n_cols = len(rows), len(rows[0])
    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(0.7), Inches(1.6),
                                         Inches(12.0), Inches(3.4))
    table = table_shape.table
    table.columns[0].width = Inches(4.2)
    for c in range(1, n_cols):
        table.columns[c].width = Inches(2.6)
    for r in range(n_rows):
        for c in range(n_cols):
            cell = table.cell(r, c)
            cell.text = rows[r][c]
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            run = para.runs[0]
            is_header = r == 0
            is_best = r == n_rows - 1
            _set_font(run, 13 if not is_header else 14, bold=is_header or is_best,
                      color=WHITE if is_header else (GREEN if is_best else DARK))
            if is_header:
                cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT
            elif is_best:
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xE6, 0xF4, 0xEA)
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
    add_textbox(slide, Inches(0.7), Inches(5.3), Inches(12), Inches(1.6),
                [("最终版本：全部车道线 Precision 88.2% / Recall 76.3%；黄线 P 90.0% / R 75.0%。", 16, True, DARK),
                 ("轻量误检过滤仅删除 6 条误检（FP），正确数 TP 保持 90 条不变，因此 Precision 提升而 Recall 不降。", 14, False, GREY)])


def image_grid_slides(prs, items):
    # 6 images per slide (3 cols x 2 rows)
    cols, rows_per = 3, 2
    per_slide = cols * rows_per
    margin_x, top0 = Inches(0.35), Inches(1.25)
    cell_w = Inches((13.333 - 0.7) / cols)
    cell_h = Inches((7.5 - 1.4) / rows_per)
    img_h = cell_h - Inches(0.42)
    for start in range(0, len(items), per_slide):
        chunk = items[start:start + per_slide]
        page_no = start // per_slide + 1
        total_pages = (len(items) + per_slide - 1) // per_slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        section_title(slide, "37 张检测结果（%d/%d）" % (page_no, total_pages))
        for idx, item in enumerate(chunk):
            col = idx % cols
            row = idx // cols
            cx = Emu(int(margin_x) + col * int(cell_w))
            cy = Emu(int(top0) + row * int(cell_h))
            # caption
            cap = "图%d  车道%d 白%d 黄%d" % (item["no"], item["lanes"], item["white"], item["yellow"])
            # place image, fit width
            pic = slide.shapes.add_picture(item["path"], cx + Inches(0.05), cy,
                                           width=cell_w - Inches(0.1))
            if pic.height > int(img_h):
                ratio = int(img_h) / pic.height
                pic.height = int(img_h)
                pic.width = int(pic.width * ratio)
                pic.left = Emu(int(cx) + (int(cell_w) - pic.width) // 2)
            add_textbox(slide, cx, cy + img_h + Inches(0.02), cell_w, Inches(0.38),
                        [(cap, 11, False, DARK)], align=PP_ALIGN.CENTER)


def main():
    summary_keys_path = os.path.join(RESULTS_DIR, "results.json")
    with open(summary_keys_path, encoding="utf-8") as f:
        data = json.load(f)
    keys = [k for k in data if k != "__summary__"]
    items = []
    for i, k in enumerate(keys, 1):
        vis = os.path.join(RESULTS_DIR, os.path.splitext(k)[0] + "_vis.jpg")
        items.append({"no": i, "path": vis, "lanes": data[k]["num_lanes"],
                      "white": data[k]["num_white"], "yellow": data[k]["num_yellow"]})

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    title_slide(prs)

    bullets_slide(prs, "任务背景与要求", [
        ("以司机视角路面图片为输入，检测车道线并判断颜色（仅黄/白两类）", 0, 20, True, DARK),
        ("可使用基于 CULane 数据集的预训练模型", 0, 18, False, GREY),
        ("对每张图输出可视化结果，标示车道线与每条线颜色", 0, 18, False, GREY),
        ("考虑运行效率：批量测试并统计平均单张处理时间", 0, 18, False, GREY),
        ("统计指标：车道线/白线/黄线的精确率(Precision)与召回率(Recall)", 0, 18, False, GREY),
        ("评判标准：线偏离较大（约 15° 以外）不计为正确", 1, 16, False, GREY),
    ])

    bullets_slide(prs, "方法与流程", [
        ("检测模型：CLRNet（CULane ResNet18），导出 ONNX，CPU 推理", 0, 19, True, DARK),
        ("输入 800×320，置信度阈值 0.30，NMS 后最多保留 8 条候选", 1, 15, False, GREY),
        ("颜色分类：道路相对蓝色亏损法", 0, 19, True, DARK),
        ("以车道线旁路面为参考做相对比较，抵消全局暖光色偏；仅取车道线中段采样", 1, 15, False, GREY),
        ("轻量误检后处理：通用几何启发式规则（非按图硬编码）", 0, 19, True, DARK),
        ("过滤贴右边界高位结束的黄线（护栏）、贴左边界乡道路缘白线、极右重复边界线（反光）", 1, 15, False, GREY),
        ("可视化：原图绘制每条线 + 标注 yellow/white；批量输出 results.json 与耗时", 0, 18, False, GREY),
    ])

    metrics_table_slide(prs)

    image_grid_slides(prs, items)

    bullets_slide(prs, "局限与诚实性说明", [
        ("轻量误检过滤是通用几何规则，代码中无任何按图片硬编码删除", 0, 19, True, DARK),
        ("但其数字阈值是在本 test-difficult 测试集上调参得到的", 0, 18, False, GREY),
        ("→ 存在过拟合到该测试集的风险：换分布不同的新图不保证同样有效，甚至可能误删真实车道线", 1, 16, False, RGBColor(0xB0, 0x47, 0x00)),
        ("它本质是后处理“打补丁”，而非检测模型本身变强", 0, 18, False, GREY),
        ("残留问题：部分乡道路缘/护栏仍有少量误检；强光、大车遮挡场景仍会漏检外侧/远处车道", 0, 18, False, GREY),
        ("更稳健的方向：更强检测模型或在更大数据上重训/校准", 0, 18, False, GREY),
    ])

    bullets_slide(prs, "结论", [
        ("CLRNet 替换 UFLD 后多车道召回明显提升，颜色分类与轻量过滤进一步提高精确率", 0, 19, True, DARK),
        ("最终（test-difficult 37 张，对照老师真值）：", 0, 18, False, GREY),
        ("全部车道线 P / R = 88.2% / 76.3%", 1, 18, True, GREEN),
        ("黄线 P / R = 90.0% / 75.0%；白线 P / R = 88.0% / 76.4%", 1, 16, False, GREY),
        ("平均单张耗时约 54.4 ms（约 18 FPS），满足批量测试效率要求", 0, 18, False, GREY),
    ])

    prs.save(OUT_PATH)
    print("Saved", OUT_PATH, "slides:", len(prs.slides._sldIdLst))


if __name__ == "__main__":
    main()
